"""
AI Analysis Pipeline for MediaCMS - WAIC Platform
Handles document/image/audio/video analysis via DashScope Qwen models.

Architecture:
    upload → media_init() → ai_analyze_media.delay(token) → analyze_media(media)
                                                                  │
                                          ┌───────────────────────┼───────────────────────┐
                                          ▼                       ▼                       ▼
                                    analyze_document         analyze_image         analyze_media_file
                                    (Qwen TURBO)            (Qwen VL Plus)         (audio/video)
"""

import base64
import json
import logging
import os
import tempfile
from io import BytesIO
from typing import Optional

import dashscope
import requests
from django.conf import settings
from django.core.files.base import ContentFile

logger = logging.getLogger(__name__)

# ── Document text extraction ────────────────────────────────────────────

def _extract_document_text(media) -> str:
    """Extract readable text from a document file for AI analysis.
    
    Handles: PDF, DOCX, PPTX, XLSX, TXT/MD/CSV/JSON/XML.
    Returns up to 6000 chars (enough for Qwen to summarize).
    """
    path = media.media_file.path
    ext = os.path.splitext(path)[1].lower()
    
    try:
        if ext == '.pdf':
            import fitz  # PyMuPDF
            doc = fitz.open(path)
            text = ''
            for page in doc:
                text += page.get_text()
                if len(text) > 8000:
                    break
            doc.close()
            return text[:6000]
        
        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(path)
            text = '\n'.join(p.text for p in doc.paragraphs if p.text)
            return text[:6000]
        
        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(path)
            text = '\n'.join(
                shape.text for slide in prs.slides 
                for shape in slide.shapes if hasattr(shape, 'text') and shape.text
            )
            return text[:6000]
        
        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True)
            text = ''
            for sheet_name in wb.sheetnames[:3]:  # first 3 sheets
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    row_text = ' | '.join(str(c) for c in row if c is not None)
                    if row_text:
                        text += row_text + '\n'
                    if len(text) > 8000:
                        break
                if len(text) > 8000:
                    break
            wb.close()
            return text[:6000]
        
        elif ext in ('.txt', '.md', '.csv', '.json', '.xml', '.log'):
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read(6000)
        
    except Exception as e:
        logger.warning(f"Text extraction failed for {path}: {e}")
    
    return ''


# ── Metadata extraction (no AI needed) ──────────────────────────────────

def _extract_file_metadata(media) -> dict:
    """Extract basic file metadata without AI."""
    path = media.media_file.path
    ext = os.path.splitext(path)[1].lower()
    meta = {}
    
    try:
        if ext == '.pdf':
            import fitz
            doc = fitz.open(path)
            meta['page_count'] = doc.page_count
            meta['author'] = doc.metadata.get('author', '')
            doc.close()
        
        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(path)
            meta['page_count'] = len(doc.paragraphs)
            core = doc.core_properties
            meta['author'] = str(core.author) if core.author else ''
        
        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(path)
            meta['page_count'] = len(prs.slides)
        
        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True)
            meta['sheet_count'] = len(wb.sheetnames)
            wb.close()
    
    except Exception as e:
        logger.warning(f"Metadata extraction failed: {e}")
    
    # File size
    try:
        meta['file_size_bytes'] = media.media_file.size
    except Exception:
        pass
    
    return meta


# ── DashScope API helpers ────────────────────────────────────────────────

def _get_api_key():
    """Get DashScope API key from settings or env."""
    return getattr(settings, 'DASHSCOPE_API_KEY', os.environ.get('DASHSCOPE_API_KEY', ''))


def _call_qwen_text(prompt: str, max_tokens: int = 500) -> Optional[str]:
    """Call Qwen TURBO for text analysis."""
    api_key = _get_api_key()
    if not api_key:
        logger.error("DASHSCOPE_API_KEY not configured")
        return None
    
    try:
        resp = dashscope.Generation.call(
            model='qwen-plus',
            api_key=api_key,
            prompt=prompt,
            result_format='message',
            max_tokens=max_tokens,
            temperature=0.3,
        )
        if resp.status_code == 200:
            return resp.output.choices[0].message.content
        else:
            logger.error(f"Qwen API error: {resp.code} - {resp.message}")
            return None
    except Exception as e:
        logger.error(f"Qwen call failed: {e}")
        return None


def _call_qwen_vision(image_path: str, prompt: str) -> Optional[str]:
    """Call Qwen VL Plus for image analysis using local file."""
    api_key = _get_api_key()
    if not api_key:
        logger.error("DASHSCOPE_API_KEY not configured")
        return None
    
    try:
        # Read and encode image
        with open(image_path, 'rb') as f:
            image_data = base64.b64encode(f.read()).decode('utf-8')
        
        messages = [{
            'role': 'user',
            'content': [
                {'image': f'data:image/jpeg;base64,{image_data}'},
                {'text': prompt},
            ]
        }]
        
        resp = dashscope.MultiModalConversation.call(
            model='qwen-vl-plus',
            api_key=api_key,
            messages=messages,
            max_tokens=500,
        )
        if resp.status_code == 200:
            return resp.output.choices[0].message.content[0]['text']
        else:
            logger.error(f"Qwen VL error: {resp.code} - {resp.message}")
            return None
    except Exception as e:
        logger.error(f"Qwen VL call failed: {e}")
        return None


# ── Type-specific analyzers ──────────────────────────────────────────────

def analyze_document(media) -> bool:
    """
    Analyze document-type media (PDF/DOCX/PPTX/XLSX/TXT).
    Extract text → AI with existing tags/categories → match/create tags + categorize.
    """
    text = _extract_document_text(media)
    file_meta = _extract_file_metadata(media)
    
    if not text:
        logger.warning(f"No text extracted from {media.friendly_token}")
        return False
    
    # Gather existing tags and categories for AI context
    tags_text, cats_text, tag_name_set, category_map = _get_existing_context()
    
    prompt = _build_tag_category_prompt(text, tags_text, cats_text, content_type="文档")
    
    result = _call_qwen_text(prompt, max_tokens=600)
    if not result:
        return False
    
    # Parse JSON response
    try:
        result = result.strip()
        if result.startswith('```'):
            result = result.split('\n', 1)[1]
            if result.endswith('```'):
                result = result[:-3]
        data = json.loads(result)
    except json.JSONDecodeError:
        logger.error(f"Failed to parse Qwen response: {result[:200]}")
        return False
    
    # Write to model
    media.ai_summary = data.get('summary', '')[:500]
    media.language = data.get('language', '')[:10]
    
    # Build ai_metadata (without ai_tags — we use real Tag model now)
    ai_meta = dict(file_meta)
    ai_meta['keywords'] = data.get('keywords', [])
    ai_meta['category_reason'] = data.get('category_reason', '')
    media.ai_metadata = ai_meta
    
    # Apply tags (match-first) and category (match or create)
    _auto_tag(media, data)
    _auto_categorize(
        media,
        category_match=data.get('category_match'),
        category_suggestion=data.get('category_suggestion'),
    )
    
    media.save(update_fields=['ai_summary', 'language', 'ai_metadata'])
    return True


def analyze_image(media) -> bool:
    """
    Analyze image-type media using Qwen VL Plus.
    AI sees both the image AND existing tags/categories → match/create tags + categorize.
    """
    # Gather existing tags and categories
    tags_text, cats_text, tag_name_set, category_map = _get_existing_context()
    
    prompt = f"""你是一个中文内容分类助手。请分析这张图片，为它分配合适的标签和分类。

== 系统中已有的标签 ==
{tags_text}

== 系统中已有的分类 ==
{cats_text}

== 你的任务 ==
返回一个 JSON 对象（只返回 JSON，不要其他文字）：

{{
  "description": "图片场景描述（中文，不超过150字）",
  "objects": ["物体1", "物体2"],
  "scene_type": "indoor/outdoor/document/screenshot/chart/portrait/landscape/abstract",
  "dominant_colors": ["#RRGGBB"],
  "tags_to_apply": ["标签1", "标签2"],
  "tags_to_create": ["新标签名"],
  "tags_to_skip": ["模糊标签"],
  "category_match": "已有分类名或null",
  "category_suggestion": "建议新建的分类名或null",
  "category_reason": "分类依据的一句话说明"
}}

规则：
1. tags_to_apply: 从「已有标签」中选最匹配的 2-5 个，名称必须完全一致
2. tags_to_create: 需要但找不到的，建议新建 1-3 个
3. category_match: 最合适的已有分类名或 null
4. 所有标签和分类名称必须使用中文"""
    
    result = _call_qwen_vision(media.media_file.path, prompt)
    if not result:
        return False
    
    try:
        result = result.strip()
        if result.startswith('```'):
            result = result.split('\n', 1)[1]
            if result.endswith('```'):
                result = result[:-3]
        data = json.loads(result)
    except json.JSONDecodeError:
        logger.error(f"Failed to parse VL response: {result[:200]}")
        return False
    
    media.ai_summary = data.get('description', '')[:500]
    
    ai_meta = {
        'objects': data.get('objects', []),
        'scene_type': data.get('scene_type', ''),
        'dominant_colors': data.get('dominant_colors', []),
        'category_reason': data.get('category_reason', ''),
    }
    media.ai_metadata = ai_meta
    
    _auto_tag(media, data)
    _auto_categorize(
        media,
        category_match=data.get('category_match'),
        category_suggestion=data.get('category_suggestion'),
    )
    
    media.save(update_fields=['ai_summary', 'ai_metadata'])
    return True


def analyze_media_file(media) -> bool:
    """
    Analyze audio/video media.
    For now: generate a content-aware summary based on filename + metadata.
    Full transcription (Paraformer) to be added in Phase 1.5.
    """
    # Use filename + media_info for basic analysis
    title = media.title or os.path.basename(media.media_file.name)
    
    prompt = f"""Analyze this media file based on its filename and suggest content metadata.
Filename: {title}
Media type: {media.media_type}
Duration: {media.duration}s (0 means unknown)

Return a JSON object with:
- summary: A 1-2 sentence guess about the content (use "likely contains" language)
- language: Most probable language (zh/en/ja)
- tags: Array of 3-5 keyword tags

Return ONLY valid JSON, no other text."""

    result = _call_qwen_text(prompt, max_tokens=300)
    if not result:
        return False
    
    try:
        result = result.strip()
        if result.startswith('```'):
            result = result.split('\n', 1)[1]
            if result.endswith('```'):
                result = result[:-3]
        data = json.loads(result)
    except json.JSONDecodeError:
        logger.error(f"Failed to parse media analysis: {result[:200]}")
        return False
    
    media.ai_summary = data.get('summary', '')[:500]
    media.language = data.get('language', '')[:10]
    
    ai_meta = {'ai_tags': data.get('tags', []), 'duration': media.duration}
    media.ai_metadata = ai_meta
    
    _auto_tag(media, data.get('tags', []))
    
    media.save(update_fields=['ai_summary', 'language', 'ai_metadata'])
    return True


# ── Context gathering for AI ──────────────────────────────────────────────

def _get_existing_context() -> tuple:
    """Gather existing tags and categories as context for AI decisions.
    
    Returns (tags_text, categories_text, tag_name_set, category_map).
    The AI uses this to match vs create new tags/categories.
    """
    from files.models import Tag, Category
    
    # Existing tags (top 100 by usage)
    tags = list(Tag.objects.all().order_by('-media_count')[:100])
    tag_names = [t.title for t in tags if t.title]
    tag_name_set = set(tag_names)
    tags_text = '\n'.join(f"  - {t}" for t in tag_names) if tag_names else "  （暂无已有标签）"
    
    # Existing categories (global ones)
    cats = list(Category.objects.filter(is_global=True))
    cat_names = [c.title for c in cats if c.title]
    cats_text = '\n'.join(f"  - {t}" for t in cat_names) if cat_names else "  （暂无已有分类）"
    category_map = {c.title: c for c in cats}
    
    return tags_text, cats_text, tag_name_set, category_map


# ── Tag management ───────────────────────────────────────────────────────

def _auto_tag(media, ai_tag_data: dict) -> None:
    """Apply AI tag decisions with match-first logic.
    
    ai_tag_data keys:
    - tags_to_apply: list[str] — existing tag names AI matched
    - tags_to_create: list[str] — new tag names AI recommends creating
    - tags_to_skip: list[str] — tags AI considered but rejected (logged only)
    """
    from files.models import Tag
    
    applied_count = 0
    created_count = 0
    
    # 1. Apply matched existing tags
    for name in (ai_tag_data.get('tags_to_apply') or []):
        name = name.strip()[:50]
        if not name:
            continue
        tag = Tag.objects.filter(title=name).first()
        if tag:
            media.tags.add(tag)
            applied_count += 1
    
    # 2. Create genuinely new tags
    for name in (ai_tag_data.get('tags_to_create') or []):
        name = name.strip()[:50]
        if not name:
            continue
        tag, created = Tag.objects.get_or_create(
            title=name,
            defaults={'source': 'ai', 'user': media.user}
        )
        if created:
            confidence = min(0.95, 0.65 + len(name) * 0.02)
            Tag.objects.filter(pk=tag.pk).update(source='ai', confidence=confidence)
            created_count += 1
        media.tags.add(tag)
    
    # 3. Log skipped
    skipped = ai_tag_data.get('tags_to_skip') or []
    if skipped:
        logger.info(f"AI skipped tags for {media.friendly_token}: {skipped}")
    
    logger.info(f"Auto-tag {media.friendly_token}: {applied_count} matched, {created_count} created, {len(skipped)} skipped")


# ── Auto-categorization ──────────────────────────────────────────────────

def _auto_categorize(media, category_match: str = '', category_suggestion: str = '') -> None:
    """Apply AI category decision.
    
    category_match: exact title of an existing Category in DB
    category_suggestion: new category name to create if no match found
    """
    from files.models import Category
    
    if category_match:
        cat = Category.objects.filter(title=category_match).first()
        if cat:
            media.category.add(cat)
            logger.info(f"Auto-categorized {media.friendly_token} → existing '{category_match}'")
            return
    
    # No match → create new category if AI suggested one
    if category_suggestion:
        cat, created = Category.objects.get_or_create(
            title=category_suggestion,
            defaults={'is_global': True}
        )
        media.category.add(cat)
        if created:
            logger.info(f"Auto-created category '{category_suggestion}' for {media.friendly_token}")
        else:
            logger.info(f"Auto-categorized {media.friendly_token} → '{category_suggestion}'")


# ── Unified AI decision prompt builder ───────────────────────────────────

def _build_tag_category_prompt(content_text: str, existing_tags: str, existing_cats: str,
                                content_type: str = "document") -> str:
    """Build the AI prompt for tag/category decisions."""
    return f"""你是一个中文内容分类助手。请分析以下{content_type}内容，为它分配合适的标签和分类。

== 系统中已有的标签 ==
{existing_tags}

== 系统中已有的分类 ==
{existing_cats}

== 需要分析的内容 ==
{content_text[:4000]}

== 你的任务 ==
返回一个 JSON 对象（只返回 JSON，不要其他文字）：

{{
  "summary": "2-3句话的中文摘要（不超过200字）",
  "language": "内容语言代码（zh/en/ja等）",
  "tags_to_apply": ["标签1", "标签2"],
  "tags_to_create": ["新标签名"],
  "tags_to_skip": ["模糊标签"],
  "category_match": "已有分类名或null",
  "category_suggestion": "建议新建的分类名或null",
  "category_reason": "分类依据的一句话说明"
}}

规则：
1. tags_to_apply: 从「已有标签」中选出最匹配的 2-5 个，名称必须完全一致
2. tags_to_create: 内容明显需要、但已有标签中找不到的，建议新建 1-3 个（中文、简洁）
3. tags_to_skip: AI 想到了但觉得不够好/过于宽泛的标签（可选）
4. category_match: 从中选一个最合适的分类，名称必须完全一致；如果没有则填 null
5. category_suggestion: 如果 category_match 为 null，建议一个新建的分类名（中文、3-8字）
6. 所有标签和分类名称必须使用中文"""


# ── Type-specific analyzers ──────────────────────────────────────────────


# ── Main entry point ─────────────────────────────────────────────────────

def analyze_media(media) -> bool:
    """
    Main AI analysis dispatcher. Called from Celery task.
    Routes to the correct analyzer based on media_type.
    """
    logger.info(f"AI analysis started for {media.friendly_token} (type={media.media_type})")
    
    try:
        if media.media_type in ('document', 'text', 'spreadsheet', 'presentation', 'pdf'):
            success = analyze_document(media)
        elif media.media_type == 'image':
            success = analyze_image(media)
        elif media.media_type in ('audio', 'video'):
            success = analyze_media_file(media)
        else:
            logger.warning(f"Unknown media_type '{media.media_type}' for {media.friendly_token}")
            return False
        
        if success:
            # Tags and categories already applied inside each analyzer.
            # Update search vector to include AI fields
            media.update_search_vector()

            # Generate embedding for semantic search
            from files.tasks import embed_media as _embed_media
            _embed_media.apply_async(args=[media.friendly_token], countdown=5)

            logger.info(f"AI analysis completed for {media.friendly_token}")
        else:
            logger.warning(f"AI analysis returned no results for {media.friendly_token}")

        return success

    except Exception as e:
        logger.exception(f"AI analysis failed for {media.friendly_token}: {e}")
        return False


# ── Embedding generation ─────────────────────────────────────────────────

def generate_embedding(text: str) -> Optional[list]:
    """Generate 1536-dim text embedding via DashScope text-embedding-v3.

    Returns a list of 1536 floats, or None on failure.
    """
    api_key = _get_api_key()
    if not api_key:
        logger.error("Cannot generate embedding: DASHSCOPE_API_KEY not configured")
        return None

    import http.client
    import socket

    try:
        body = json.dumps({
            "model": "text-embedding-v3",
            "input": {"texts": [text[:6000]]},
            "parameters": {"text_type": "document"}
        })
        old_timeout = socket.getdefaulttimeout()
        socket.setdefaulttimeout(8)
        try:
            conn = http.client.HTTPSConnection("dashscope.aliyuncs.com", timeout=5)
            conn.request(
                "POST",
                "/api/v1/services/embeddings/text-embedding/text-embedding",
                body=body,
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
            )
            resp = conn.getresponse()
            result = json.loads(resp.read().decode("utf-8"))
            conn.close()
        finally:
            socket.setdefaulttimeout(old_timeout)

        if "output" in result and result["output"].get("embeddings"):
            return result["output"]["embeddings"][0]["embedding"]
        else:
            logger.error(f"Embedding API error: {result.get('message', 'Unknown')}")
            return None
    except Exception as e:
        logger.error(f"Embedding generation failed: {e}")
        return None


def embed_media(media) -> bool:
    """Generate and save embedding for a media item.

    Builds text from title + ai_summary + tags, then calls DashScope.
    Returns True on success.
    """
    parts = [media.title or ""]
    if media.ai_summary:
        parts.append(media.ai_summary)
    # Use real tags from ManyToMany relationship
    tag_names = [t.title for t in media.tags.all()[:10]]
    if tag_names:
        parts.append(" ".join(tag_names))

    text = " ".join(p for p in parts if p)
    if not text:
        logger.warning(f"No text to embed for {media.friendly_token}")
        return False

    embedding = generate_embedding(text)
    if embedding:
        # Use update() to avoid triggering signals
        from files.models import Media
        Media.objects.filter(pk=media.pk).update(embedding=embedding)
        logger.info(f"Embedding saved for {media.friendly_token}")
        return True
    return False
